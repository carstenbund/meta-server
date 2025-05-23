#!/usr/bin/env python3

import os
import time
import threading
import logging
import requests
from sqlalchemy import create_engine, select
from sqlalchemy.orm import sessionmaker, declarative_base
from sqlalchemy.exc import SQLAlchemyError, IntegrityError
from langdetect import detect
from langdetect.lang_detect_exception import LangDetectException
from tika import parser
import magic
import pefile
import openpyxl
from pptx import Presentation
from MyLogger import Logger

# ---------- Configuration ----------
DATABASE_URI = 'sqlite:///instance/files.db'
INFERENCE_URL = 'http://localhost:5001/infer'
NUM_WORKERS = 4
POLL_INTERVAL = 2  # seconds

# ---------- Logger Setup ----------
log = Logger(log_name='index_worker', log_level=logging.DEBUG).get_logger()

# ---------- SQLAlchemy Setup ----------
Base = declarative_base()
engine = create_engine(
    DATABASE_URI,
    connect_args={"check_same_thread": False},
    pool_pre_ping=True
)
Session = sessionmaker(bind=engine)

# ---------- Models ----------
class FileMetadata(Base):
    __tablename__ = 'file_metadata'
    from sqlalchemy import Column, Integer, String, Float, Text
    id = Column(Integer, primary_key=True, autoincrement=True)
    path = Column(String, unique=True, nullable=False)
    size = Column(Integer, nullable=False)
    modification_date = Column(Float, nullable=False)
    category = Column(String)
    inferred_category = Column(String)
    keywords = Column(String)
    summary = Column(String)
    content = Column(Text)
    file_type = Column(String)
    creator_software = Column(String)
    origin_date = Column(String)
    pe_info = Column(Text)

class IndexQueue(Base):
    __tablename__ = 'index_queue'
    from sqlalchemy import Column, Integer, String, Float, Text
    id = Column(Integer, primary_key=True)
    file_path = Column(String, unique=True, nullable=False)
    status = Column(String, default='pending')  # pending, in_progress, done, error
    error = Column(Text, nullable=True)
    added_at = Column(Float)
    started_at = Column(Float)
    finished_at = Column(Float)

# ---------- Helper Functions ----------
def detect_file_type(file_path):
    file_magic = magic.Magic(mime=True)
    mime_type = file_magic.from_file(file_path)
    parts = mime_type.split('/')
    file_type = parts[0]
    creator_software = parts[1] if len(parts) > 1 else 'Unknown'
    return mime_type, file_type, creator_software

def get_pe_info(file_path):
    try:
        pe = pefile.PE(file_path)
        return str({
            "entry_point": pe.OPTIONAL_HEADER.AddressOfEntryPoint,
            "image_base": pe.OPTIONAL_HEADER.ImageBase,
            "number_of_sections": pe.FILE_HEADER.NumberOfSections
        })
    except Exception as e:
        return str(e)

def clean_text(text):
    import re
    if not isinstance(text, str):
        return ""
    text = re.sub(r'\n+', '\n', text)
    text = text.strip()
    return re.sub(r'\s+', ' ', text)

# ---------- Indexing Thread ----------
class IndexWorker(threading.Thread):
    def __init__(self, worker_id):
        super().__init__()
        self.worker_id = worker_id
        self.running = True

    def run(self):
        log.info(f"Worker-{self.worker_id} started.")
        while self.running:
            with Session() as session:
                job = session.execute(
                    select(IndexQueue)
                    .where(IndexQueue.status == 'pending')
                    .limit(1)
                ).scalar_one_or_none()

                if not job:
                    time.sleep(POLL_INTERVAL)
                    continue

                log.debug(f"Worker-{self.worker_id} processing {job.file_path}")
                job.status = 'in_progress'
                job.started_at = time.time()
                session.commit()

                try:
                    self.process_file(session, job)
                    job.status = 'done'
                    job.finished_at = time.time()
                    session.commit()
                    log.info(f"Worker-{self.worker_id} indexed {job.file_path}")
                except Exception as e:
                    session.rollback()
                    job.status = 'error'
                    job.error = str(e)
                    job.finished_at = time.time()
                    session.commit()
                    log.exception(f"Worker-{self.worker_id} failed on {job.file_path}: {e}")



    def process_file(self, session, job):
        path = job.file_path
        if not os.path.exists(path):
            raise FileNotFoundError(f"File does not exist: {path}")

        size = os.path.getsize(path)
        mod_time = os.path.getmtime(path)
        origin_date = str(mod_time)
        mime_type, file_type, creator_software = detect_file_type(path)
        category = os.path.basename(os.path.dirname(path))
        pe_info = ""
        content = ""
        inferred_category = None
        keywords = None
        summary = None

        ext = os.path.splitext(path)[1].lower()

        if ext == ".exe":
            pe_info = get_pe_info(path)
            content = pe_info
            inferred_category = "Executable"

        elif ext in [".pdf", ".doc", ".docx"]:
            parsed = parser.from_file(path)
            raw_content = parsed.get("content", "") if parsed else ""
            content = clean_text(raw_content)

            try:
                lang = detect(content[:500])
            except LangDetectException:
                lang = "unknown"

            payload = {
                "file_path": path,
                "content": content[:500],
                "language": lang
            }

            response = requests.post(INFERENCE_URL, json=payload)
            if response.status_code == 200:
                result = response.json()
                inferred_category = result.get("category")
                keywords = result.get("keywords")
                summary = result.get("summary")
                if summary:
                    content = summary
            else:
                raise RuntimeError(f"Inference failed: {response.status_code} {response.text}")

        elif ext == ".pptx":
            from pptx import Presentation
            content = ""
            try:
                prs = Presentation(path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                content = clean_text(content)
            except Exception as e:
                log.warning(f"pptx parsing failed for {path}: {e}")
                content = ""

        elif ext == ".xlsx":
            import openpyxl
            content = ""
            try:
                wb = openpyxl.load_workbook(path, read_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        line = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        content += line + "\n"
                content = clean_text(content)
            except Exception as e:
                log.warning(f"xlsx parsing failed for {path}: {e}")
                content = ""

        elif ext in [".txt", ".md"]:
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = clean_text(f.read())
            except Exception as e:
                log.warning(f"txt/md parsing failed for {path}: {e}")
                content = ""

        else:
            # fallback to tika
            parsed = parser.from_file(path)
            raw_content = parsed.get("content", "") if parsed else ""
            content = clean_text(raw_content)

        # Safe upsert
        existing = session.query(FileMetadata).filter_by(path=path).first()
        if existing:
            for key, value in {
                "size": size,
                "modification_date": mod_time,
                "category": category,
                "inferred_category": inferred_category,
                "keywords": keywords,
                "summary": summary,
                "content": content,
                "file_type": file_type,
                "creator_software": creator_software,
                "origin_date": origin_date,
                "pe_info": pe_info,
            }.items():
                setattr(existing, key, value)
        else:
            metadata = FileMetadata(
                path=path,
                size=size,
                modification_date=mod_time,
                category=category,
                inferred_category=inferred_category,
                keywords=keywords,
                summary=summary,
                content=content,
                file_type=file_type,
                creator_software=creator_software,
                origin_date=origin_date,
                pe_info=pe_info
            )
            session.add(metadata)


# ---------- Entrypoint ----------
if __name__ == '__main__':
    Base.metadata.create_all(engine)
    log.info("Indexing service starting...")

    workers = [IndexWorker(i) for i in range(NUM_WORKERS)]
    for w in workers:
        w.start()

    try:
        while True:
            time.sleep(5)
    except KeyboardInterrupt:
        log.info("Stopping all workers...")
        for w in workers:
            w.running = False
        for w in workers:
            w.join()
        log.info("Indexing service stopped.")

